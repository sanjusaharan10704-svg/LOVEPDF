"""
Comprehensive backend test for PDF tools API endpoints.
Tests all endpoints under /api/pdf/* with real file generation.
"""
import os
import sys
import requests
import tempfile
from pathlib import Path
from io import BytesIO

# Backend URL from frontend .env
BACKEND_URL = "https://document-master-13.preview.emergentagent.com"
BASE_URL = f"{BACKEND_URL}/api/pdf"

# Test results tracking
test_results = {
    "passed": [],
    "failed": [],
    "warnings": []
}

def log_pass(test_name, details=""):
    print(f"✅ PASS: {test_name}")
    if details:
        print(f"   {details}")
    test_results["passed"].append(test_name)

def log_fail(test_name, error):
    print(f"❌ FAIL: {test_name}")
    print(f"   Error: {error}")
    test_results["failed"].append(f"{test_name}: {error}")

def log_warning(test_name, warning):
    print(f"⚠️  WARNING: {test_name}")
    print(f"   {warning}")
    test_results["warnings"].append(f"{test_name}: {warning}")

def check_pdf_signature(data):
    """Check if data starts with PDF magic bytes"""
    return data[:4] == b'%PDF'

def check_zip_signature(data):
    """Check if data starts with ZIP magic bytes (for docx/xlsx/pptx)"""
    return data[:2] == b'PK'

def generate_test_docx():
    """Generate a small test .docx file"""
    from docx import Document
    doc = Document()
    doc.add_heading('Test Document for PDFPro', 0)
    doc.add_paragraph('This is a test document created for testing office-to-pdf conversion.')
    doc.add_paragraph('It contains multiple paragraphs to ensure proper conversion.')
    
    temp_path = Path(tempfile.gettempdir()) / "test_document.docx"
    doc.save(str(temp_path))
    return temp_path

def generate_test_xlsx():
    """Generate a small test .xlsx file"""
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.title = "Test Sheet"
    ws['A1'] = "Name"
    ws['B1'] = "Value"
    ws['A2'] = "Test Item 1"
    ws['B2'] = 100
    ws['A3'] = "Test Item 2"
    ws['B3'] = 200
    
    temp_path = Path(tempfile.gettempdir()) / "test_spreadsheet.xlsx"
    wb.save(str(temp_path))
    return temp_path

def generate_test_pptx():
    """Generate a small test .pptx file"""
    from pptx import Presentation
    from pptx.util import Inches
    
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Test Presentation"
    subtitle.text = "Created for PDFPro testing"
    
    temp_path = Path(tempfile.gettempdir()) / "test_presentation.pptx"
    prs.save(str(temp_path))
    return temp_path

def generate_test_pdf():
    """Generate a multi-page test PDF"""
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter
    
    temp_path = Path(tempfile.gettempdir()) / "test_document.pdf"
    c = canvas.Canvas(str(temp_path), pagesize=letter)
    
    # Page 1
    c.drawString(100, 750, "PDFPro Test Document - Page 1")
    c.drawString(100, 700, "This is a multi-page PDF for testing.")
    c.drawString(100, 650, "Name: John Doe")
    c.drawString(100, 600, "Email: john@example.com")
    c.showPage()
    
    # Page 2
    c.drawString(100, 750, "PDFPro Test Document - Page 2")
    c.drawString(100, 700, "This page contains additional test data.")
    c.drawString(100, 650, "Product: Widget A")
    c.drawString(100, 600, "Price: $99.99")
    c.showPage()
    
    c.save()
    return temp_path

def generate_password_protected_pdf(password="secret123"):
    """Generate a password-protected PDF using pikepdf"""
    import pikepdf
    
    # First create a simple PDF
    source_pdf = generate_test_pdf()
    
    # Now encrypt it
    temp_path = Path(tempfile.gettempdir()) / "test_protected.pdf"
    with pikepdf.open(str(source_pdf)) as pdf:
        pdf.save(str(temp_path), encryption=pikepdf.Encryption(owner=password, user=password, R=6))
    
    return temp_path

def generate_test_image():
    """Generate a simple test image"""
    from PIL import Image, ImageDraw, ImageFont
    
    img = Image.new('RGB', (400, 300), color='white')
    draw = ImageDraw.Draw(img)
    draw.rectangle([50, 50, 350, 250], outline='black', width=2)
    draw.text((100, 130), "Test Image for OCR", fill='black')
    
    temp_path = Path(tempfile.gettempdir()) / "test_image.png"
    img.save(str(temp_path))
    return temp_path

print("=" * 80)
print("PDF TOOLS BACKEND API TEST SUITE")
print("=" * 80)
print(f"Backend URL: {BACKEND_URL}")
print(f"API Base: {BASE_URL}")
print("=" * 80)

# Generate test files
print("\n📁 Generating test files...")
try:
    docx_file = generate_test_docx()
    print(f"   ✓ Created: {docx_file}")
    xlsx_file = generate_test_xlsx()
    print(f"   ✓ Created: {xlsx_file}")
    pptx_file = generate_test_pptx()
    print(f"   ✓ Created: {pptx_file}")
    pdf_file = generate_test_pdf()
    print(f"   ✓ Created: {pdf_file}")
    protected_pdf = generate_password_protected_pdf()
    print(f"   ✓ Created: {protected_pdf}")
    image_file = generate_test_image()
    print(f"   ✓ Created: {image_file}")
except Exception as e:
    print(f"❌ Failed to generate test files: {e}")
    sys.exit(1)

print("\n" + "=" * 80)
print("RUNNING TESTS")
print("=" * 80)

# Test 1: Health Check
print("\n[1/13] Testing GET /api/pdf/health")
try:
    response = requests.get(f"{BASE_URL}/health", timeout=30)
    if response.status_code == 200:
        data = response.json()
        if data.get("ok") and "tools" in data:
            tools = data["tools"]
            all_tools_available = all([
                tools.get("soffice"),
                tools.get("gs"),
                tools.get("qpdf"),
                tools.get("tesseract"),
                tools.get("pdftoppm"),
                tools.get("ocrmypdf")
            ])
            if all_tools_available:
                log_pass("Health check", f"All tools available: {tools}")
            else:
                log_warning("Health check", f"Some tools missing: {tools}")
        else:
            log_fail("Health check", f"Invalid response format: {data}")
    else:
        log_fail("Health check", f"HTTP {response.status_code}: {response.text}")
except Exception as e:
    log_fail("Health check", str(e))

# Test 2: Office to PDF - DOCX
print("\n[2/13] Testing POST /api/pdf/office-to-pdf (DOCX)")
try:
    with open(docx_file, 'rb') as f:
        files = {'file': ('test.docx', f, 'application/vnd.openxmlformats-officedocument.wordprocessingml.document')}
        response = requests.post(f"{BASE_URL}/office-to-pdf", files=files, timeout=180)
    
    if response.status_code == 200:
        if check_pdf_signature(response.content):
            log_pass("Office to PDF (DOCX)", f"Returned valid PDF ({len(response.content)} bytes)")
        else:
            log_fail("Office to PDF (DOCX)", "Response is not a valid PDF (missing %PDF signature)")
    else:
        log_fail("Office to PDF (DOCX)", f"HTTP {response.status_code}: {response.text}")
except Exception as e:
    log_fail("Office to PDF (DOCX)", str(e))

# Test 3: Office to PDF - XLSX
print("\n[3/13] Testing POST /api/pdf/office-to-pdf (XLSX)")
try:
    with open(xlsx_file, 'rb') as f:
        files = {'file': ('test.xlsx', f, 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')}
        response = requests.post(f"{BASE_URL}/office-to-pdf", files=files, timeout=180)
    
    if response.status_code == 200:
        if check_pdf_signature(response.content):
            log_pass("Office to PDF (XLSX)", f"Returned valid PDF ({len(response.content)} bytes)")
        else:
            log_fail("Office to PDF (XLSX)", "Response is not a valid PDF")
    else:
        log_fail("Office to PDF (XLSX)", f"HTTP {response.status_code}: {response.text}")
except Exception as e:
    log_fail("Office to PDF (XLSX)", str(e))

# Test 4: Office to PDF - PPTX
print("\n[4/13] Testing POST /api/pdf/office-to-pdf (PPTX)")
try:
    with open(pptx_file, 'rb') as f:
        files = {'file': ('test.pptx', f, 'application/vnd.openxmlformats-officedocument.presentationml.presentation')}
        response = requests.post(f"{BASE_URL}/office-to-pdf", files=files, timeout=180)
    
    if response.status_code == 200:
        if check_pdf_signature(response.content):
            log_pass("Office to PDF (PPTX)", f"Returned valid PDF ({len(response.content)} bytes)")
        else:
            log_fail("Office to PDF (PPTX)", "Response is not a valid PDF")
    else:
        log_fail("Office to PDF (PPTX)", f"HTTP {response.status_code}: {response.text}")
except Exception as e:
    log_fail("Office to PDF (PPTX)", str(e))

# Test 5: HTML to PDF (with HTML content)
print("\n[5/13] Testing POST /api/pdf/html-to-pdf (HTML content)")
try:
    data = {'html': '<h1>Hello PDFPro</h1><p>This is a test HTML to PDF conversion.</p>'}
    response = requests.post(f"{BASE_URL}/html-to-pdf", data=data, timeout=180)
    
    if response.status_code == 200:
        if check_pdf_signature(response.content):
            log_pass("HTML to PDF (content)", f"Returned valid PDF ({len(response.content)} bytes)")
        else:
            log_fail("HTML to PDF (content)", "Response is not a valid PDF")
    else:
        log_fail("HTML to PDF (content)", f"HTTP {response.status_code}: {response.text}")
except Exception as e:
    log_fail("HTML to PDF (content)", str(e))

# Test 6: HTML to PDF (with URL)
print("\n[6/13] Testing POST /api/pdf/html-to-pdf (URL)")
try:
    data = {'url': 'https://example.com'}
    response = requests.post(f"{BASE_URL}/html-to-pdf", data=data, timeout=180)
    
    if response.status_code == 200:
        if check_pdf_signature(response.content):
            log_pass("HTML to PDF (URL)", f"Returned valid PDF ({len(response.content)} bytes)")
        else:
            log_fail("HTML to PDF (URL)", "Response is not a valid PDF")
    else:
        log_fail("HTML to PDF (URL)", f"HTTP {response.status_code}: {response.text}")
except Exception as e:
    log_fail("HTML to PDF (URL)", str(e))

# Test 7: PDF to Word
print("\n[7/13] Testing POST /api/pdf/pdf-to-word")
try:
    with open(pdf_file, 'rb') as f:
        files = {'file': ('test.pdf', f, 'application/pdf')}
        response = requests.post(f"{BASE_URL}/pdf-to-word", files=files, timeout=180)
    
    if response.status_code == 200:
        if check_zip_signature(response.content):
            log_pass("PDF to Word", f"Returned valid DOCX ({len(response.content)} bytes)")
        else:
            log_fail("PDF to Word", "Response is not a valid DOCX (missing PK signature)")
    else:
        log_fail("PDF to Word", f"HTTP {response.status_code}: {response.text}")
except Exception as e:
    log_fail("PDF to Word", str(e))

# Test 8: PDF to Excel
print("\n[8/13] Testing POST /api/pdf/pdf-to-excel")
try:
    with open(pdf_file, 'rb') as f:
        files = {'file': ('test.pdf', f, 'application/pdf')}
        response = requests.post(f"{BASE_URL}/pdf-to-excel", files=files, timeout=180)
    
    if response.status_code == 200:
        if check_zip_signature(response.content):
            log_pass("PDF to Excel", f"Returned valid XLSX ({len(response.content)} bytes)")
        else:
            log_fail("PDF to Excel", "Response is not a valid XLSX")
    else:
        log_fail("PDF to Excel", f"HTTP {response.status_code}: {response.text}")
except Exception as e:
    log_fail("PDF to Excel", str(e))

# Test 9: PDF to PowerPoint
print("\n[9/13] Testing POST /api/pdf/pdf-to-ppt")
try:
    with open(pdf_file, 'rb') as f:
        files = {'file': ('test.pdf', f, 'application/pdf')}
        response = requests.post(f"{BASE_URL}/pdf-to-ppt", files=files, timeout=180)
    
    if response.status_code == 200:
        if check_zip_signature(response.content):
            log_pass("PDF to PowerPoint", f"Returned valid PPTX ({len(response.content)} bytes)")
        else:
            log_fail("PDF to PowerPoint", "Response is not a valid PPTX")
    else:
        log_fail("PDF to PowerPoint", f"HTTP {response.status_code}: {response.text}")
except Exception as e:
    log_fail("PDF to PowerPoint", str(e))

# Test 10: Protect PDF
print("\n[10/13] Testing POST /api/pdf/protect")
try:
    with open(pdf_file, 'rb') as f:
        files = {'file': ('test.pdf', f, 'application/pdf')}
        data = {'password': 'secret123'}
        response = requests.post(f"{BASE_URL}/protect", files=files, data=data, timeout=180)
    
    if response.status_code == 200:
        if check_pdf_signature(response.content):
            # Save the protected PDF for unlock test
            protected_test_pdf = Path(tempfile.gettempdir()) / "test_protected_api.pdf"
            with open(protected_test_pdf, 'wb') as f:
                f.write(response.content)
            
            # Verify it's actually encrypted by trying to open without password
            import pikepdf
            try:
                with pikepdf.open(str(protected_test_pdf)) as pdf:
                    log_fail("Protect PDF", "PDF is not encrypted (opened without password)")
            except pikepdf.PasswordError:
                # Good! It's encrypted
                try:
                    with pikepdf.open(str(protected_test_pdf), password='secret123') as pdf:
                        log_pass("Protect PDF", f"PDF properly encrypted and can be opened with correct password ({len(response.content)} bytes)")
                except Exception as e:
                    log_fail("Protect PDF", f"Cannot open with correct password: {e}")
        else:
            log_fail("Protect PDF", "Response is not a valid PDF")
    else:
        log_fail("Protect PDF", f"HTTP {response.status_code}: {response.text}")
except Exception as e:
    log_fail("Protect PDF", str(e))

# Test 11: Unlock PDF (correct password)
print("\n[11/13] Testing POST /api/pdf/unlock (correct password)")
try:
    with open(protected_pdf, 'rb') as f:
        files = {'file': ('protected.pdf', f, 'application/pdf')}
        data = {'password': 'secret123'}
        response = requests.post(f"{BASE_URL}/unlock", files=files, data=data, timeout=180)
    
    if response.status_code == 200:
        if check_pdf_signature(response.content):
            log_pass("Unlock PDF (correct password)", f"Returned decrypted PDF ({len(response.content)} bytes)")
        else:
            log_fail("Unlock PDF (correct password)", "Response is not a valid PDF")
    else:
        log_fail("Unlock PDF (correct password)", f"HTTP {response.status_code}: {response.text}")
except Exception as e:
    log_fail("Unlock PDF (correct password)", str(e))

# Test 12: Unlock PDF (wrong password)
print("\n[12/13] Testing POST /api/pdf/unlock (wrong password)")
try:
    with open(protected_pdf, 'rb') as f:
        files = {'file': ('protected.pdf', f, 'application/pdf')}
        data = {'password': 'wrongpassword'}
        response = requests.post(f"{BASE_URL}/unlock", files=files, data=data, timeout=180)
    
    if response.status_code == 400:
        log_pass("Unlock PDF (wrong password)", "Correctly returned HTTP 400 for wrong password")
    else:
        log_fail("Unlock PDF (wrong password)", f"Expected HTTP 400, got {response.status_code}")
except Exception as e:
    log_fail("Unlock PDF (wrong password)", str(e))

# Test 13: OCR
print("\n[13/13] Testing POST /api/pdf/ocr")
try:
    with open(pdf_file, 'rb') as f:
        files = {'file': ('test.pdf', f, 'application/pdf')}
        data = {'lang': 'eng'}
        response = requests.post(f"{BASE_URL}/ocr", files=files, data=data, timeout=180)
    
    if response.status_code == 200:
        if check_pdf_signature(response.content):
            log_pass("OCR PDF", f"Returned valid PDF ({len(response.content)} bytes)")
        else:
            log_fail("OCR PDF", "Response is not a valid PDF")
    else:
        log_fail("OCR PDF", f"HTTP {response.status_code}: {response.text}")
except Exception as e:
    log_fail("OCR PDF", str(e))

# Test 14: Repair
print("\n[14/13] Testing POST /api/pdf/repair")
try:
    with open(pdf_file, 'rb') as f:
        files = {'file': ('test.pdf', f, 'application/pdf')}
        response = requests.post(f"{BASE_URL}/repair", files=files, timeout=180)
    
    if response.status_code == 200:
        if check_pdf_signature(response.content):
            log_pass("Repair PDF", f"Returned valid PDF ({len(response.content)} bytes)")
        else:
            log_fail("Repair PDF", "Response is not a valid PDF")
    else:
        log_fail("Repair PDF", f"HTTP {response.status_code}: {response.text}")
except Exception as e:
    log_fail("Repair PDF", str(e))

# Test 15: PDF/A
print("\n[15/13] Testing POST /api/pdf/pdfa")
try:
    with open(pdf_file, 'rb') as f:
        files = {'file': ('test.pdf', f, 'application/pdf')}
        response = requests.post(f"{BASE_URL}/pdfa", files=files, timeout=180)
    
    if response.status_code == 200:
        if check_pdf_signature(response.content):
            log_pass("PDF/A conversion", f"Returned valid PDF ({len(response.content)} bytes)")
        else:
            log_fail("PDF/A conversion", "Response is not a valid PDF")
    else:
        log_fail("PDF/A conversion", f"HTTP {response.status_code}: {response.text}")
except Exception as e:
    log_fail("PDF/A conversion", str(e))

# Test 16: Crop
print("\n[16/13] Testing POST /api/pdf/crop")
try:
    with open(pdf_file, 'rb') as f:
        files = {'file': ('test.pdf', f, 'application/pdf')}
        data = {'margin': '5'}
        response = requests.post(f"{BASE_URL}/crop", files=files, data=data, timeout=180)
    
    if response.status_code == 200:
        if check_pdf_signature(response.content):
            log_pass("Crop PDF", f"Returned valid PDF ({len(response.content)} bytes)")
        else:
            log_fail("Crop PDF", "Response is not a valid PDF")
    else:
        log_fail("Crop PDF", f"HTTP {response.status_code}: {response.text}")
except Exception as e:
    log_fail("Crop PDF", str(e))

# Test 17: Compare
print("\n[17/13] Testing POST /api/pdf/compare")
try:
    # Create two slightly different PDFs
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter
    
    pdf1_path = Path(tempfile.gettempdir()) / "compare1.pdf"
    c1 = canvas.Canvas(str(pdf1_path), pagesize=letter)
    c1.drawString(100, 750, "Document Version 1")
    c1.drawString(100, 700, "This is the first version.")
    c1.save()
    
    pdf2_path = Path(tempfile.gettempdir()) / "compare2.pdf"
    c2 = canvas.Canvas(str(pdf2_path), pagesize=letter)
    c2.drawString(100, 750, "Document Version 2")
    c2.drawString(100, 700, "This is the second version.")
    c2.save()
    
    with open(pdf1_path, 'rb') as f1, open(pdf2_path, 'rb') as f2:
        files = {
            'file1': ('doc1.pdf', f1, 'application/pdf'),
            'file2': ('doc2.pdf', f2, 'application/pdf')
        }
        response = requests.post(f"{BASE_URL}/compare", files=files, timeout=180)
    
    if response.status_code == 200:
        data = response.json()
        if 'similarity' in data and 'rows' in data:
            if isinstance(data['similarity'], (int, float)) and isinstance(data['rows'], list):
                log_pass("Compare PDF", f"Returned valid comparison (similarity: {data['similarity']}%, {len(data['rows'])} rows)")
            else:
                log_fail("Compare PDF", f"Invalid data types in response: {data}")
        else:
            log_fail("Compare PDF", f"Missing required fields in response: {data}")
    else:
        log_fail("Compare PDF", f"HTTP {response.status_code}: {response.text}")
except Exception as e:
    log_fail("Compare PDF", str(e))

# Print summary
print("\n" + "=" * 80)
print("TEST SUMMARY")
print("=" * 80)
print(f"✅ Passed: {len(test_results['passed'])}")
print(f"❌ Failed: {len(test_results['failed'])}")
print(f"⚠️  Warnings: {len(test_results['warnings'])}")

if test_results['failed']:
    print("\n❌ FAILED TESTS:")
    for failure in test_results['failed']:
        print(f"   - {failure}")

if test_results['warnings']:
    print("\n⚠️  WARNINGS:")
    for warning in test_results['warnings']:
        print(f"   - {warning}")

print("\n" + "=" * 80)

# Exit with appropriate code
sys.exit(0 if len(test_results['failed']) == 0 else 1)
